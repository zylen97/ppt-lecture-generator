"""
PPT处理模块

负责PPT文件的读取、分析和图片转换功能。
"""

import os
import time
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
from dataclasses import dataclass

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from spire.presentation import Presentation as SpirePresentation
    from spire.presentation import FileFormat, ImageFormat
    SPIRE_AVAILABLE = True
except ImportError:
    SPIRE_AVAILABLE = False

try:
    from pdf2image import convert_from_path
    import fitz  # PyMuPDF
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

from ..utils.file_utils import FileUtils
from ..utils.logger import get_logger
from ..config.constants import PPT_PROCESSING
from ..utils.ppt_converter import PPTConverter


@dataclass
class SlideInfo:
    """幻灯片信息"""
    slide_number: int
    title: str
    content: List[str]
    bullet_points: List[Dict[str, Any]]
    image_count: int
    chart_count: int
    table_count: int
    slide_type: str
    image_path: Optional[str] = None
    notes: str = ""


class PPTProcessor:
    """PPT处理器"""
    
    def __init__(self, ppt_path: str):
        """
        初始化PPT处理器
        
        Args:
            ppt_path: PPT文件路径
        """
        self.ppt_path = Path(ppt_path)
        self.temp_dir = ""
        self.slides_info: List[SlideInfo] = []
        self.logger = get_logger()
        
        # 检查文件是否存在
        if not self.ppt_path.exists():
            raise FileNotFoundError(f"PPT文件不存在: {ppt_path}")
        
        # 检查文件格式
        if not FileUtils.is_ppt_file(ppt_path):
            raise ValueError(f"不支持的文件格式: {self.ppt_path.suffix}")
        
        # 检查文件大小
        file_size = FileUtils.get_file_size(ppt_path)
        if file_size > PPT_PROCESSING['max_file_size']:
            raise ValueError(f"文件太大: {file_size / (1024*1024):.1f}MB，最大支持{PPT_PROCESSING['max_file_size'] / (1024*1024):.0f}MB")
        
        self.logger.info(f"初始化PPT处理器: {ppt_path}")
    
    def process(self) -> bool:
        """
        处理PPT文件
        
        Returns:
            bool: 是否成功
        """
        try:
            # 创建临时目录
            self.temp_dir = FileUtils.create_temp_dir("ppt_process_")
            if not self.temp_dir:
                return False
            
            # 提取文本内容
            if not self._extract_text_content():
                return False
            
            # 转换为图片
            if not self._convert_to_images():
                return False
            
            self.logger.info(f"PPT处理完成，共{len(self.slides_info)}张幻灯片")
            return True
            
        except Exception as e:
            self.logger.error(f"PPT处理失败: {e}")
            return False
    
    def _extract_text_content(self) -> bool:
        """
        提取PPT文本内容
        
        Returns:
            bool: 是否成功
        """
        if not PPTX_AVAILABLE:
            self.logger.error("python-pptx库未安装，无法提取文本内容")
            return False
        
        try:
            presentation = Presentation(self.ppt_path)
            
            for slide_idx, slide in enumerate(presentation.slides):
                slide_info = self._analyze_slide(slide, slide_idx + 1)
                self.slides_info.append(slide_info)
            
            # 检查幻灯片数量限制
            if len(self.slides_info) > PPT_PROCESSING['max_slides']:
                self.logger.warning(f"幻灯片数量超过限制: {len(self.slides_info)}")
                self.slides_info = self.slides_info[:PPT_PROCESSING['max_slides']]
            
            return True
            
        except Exception as e:
            self.logger.error(f"提取文本内容失败: {e}")
            return False
    
    def _analyze_slide(self, slide, slide_number: int) -> SlideInfo:
        """
        分析单张幻灯片
        
        Args:
            slide: 幻灯片对象
            slide_number: 幻灯片编号
            
        Returns:
            SlideInfo: 幻灯片信息
        """
        slide_info = SlideInfo(
            slide_number=slide_number,
            title="",
            content=[],
            bullet_points=[],
            image_count=0,
            chart_count=0,
            table_count=0,
            slide_type="content"
        )
        
        # 分析形状
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    # 判断是否为标题
                    if shape == slide.shapes.title or slide_info.title == "":
                        slide_info.title = text
                        # 判断幻灯片类型
                        if "第" in text and "章" in text:
                            slide_info.slide_type = "section"
                        elif "总结" in text or "小结" in text:
                            slide_info.slide_type = "summary"
                    else:
                        # 提取段落和要点
                        for para in shape.text_frame.paragraphs:
                            if para.text.strip():
                                if para.level == 0:
                                    slide_info.content.append(para.text.strip())
                                else:
                                    slide_info.bullet_points.append({
                                        'level': para.level,
                                        'text': para.text.strip()
                                    })
            
            # 统计多媒体元素
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                slide_info.image_count += 1
            elif shape.has_chart:
                slide_info.chart_count += 1
            elif shape.has_table:
                slide_info.table_count += 1
        
        # 判断标题页
        if slide_number == 1 or (slide_info.title and not slide_info.content and not slide_info.bullet_points):
            slide_info.slide_type = "title"
        
        # 提取备注
        if slide.notes_slide:
            try:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    slide_info.notes = notes_text
            except Exception:
                pass
        
        return slide_info
    
    def _convert_to_images(self) -> bool:
        """
        将PPT转换为图片
        
        Returns:
            bool: 是否成功
        """
        # 优先使用新的PPT转换器
        try:
            converter = PPTConverter(self.logger)
            image_paths = converter.convert_ppt_to_images(
                str(self.ppt_path), 
                self.temp_dir,
                dpi=PPT_PROCESSING.get('dpi', 300)
            )
            
            if image_paths:
                # 将图片路径关联到幻灯片信息
                for i, image_path in enumerate(image_paths):
                    if i < len(self.slides_info):
                        self.slides_info[i].image_path = image_path
                return True
            else:
                self.logger.warning("PPT转换器未生成任何图片")
        except Exception as e:
            self.logger.error(f"使用PPT转换器失败: {e}")
        
        # 备选方案1：使用Spire.Presentation
        if SPIRE_AVAILABLE:
            return self._convert_with_spire()
        
        # 备选方案2：通过PDF转换
        if PDF_AVAILABLE:
            return self._convert_via_pdf()
        
        self.logger.error("没有可用的PPT转图片库")
        return False
    
    def _convert_with_spire(self) -> bool:
        """
        使用Spire.Presentation转换
        
        Returns:
            bool: 是否成功
        """
        try:
            presentation = SpirePresentation()
            presentation.LoadFromFile(str(self.ppt_path))
            
            # 设置高质量转换参数
            presentation.SlideSize.Type = presentation.SlideSize.Type
            
            for i, slide_info in enumerate(self.slides_info):
                image_path = os.path.join(self.temp_dir, f"slide_{i+1:03d}.png")
                
                # 使用高分辨率转换为图片
                # 设置缩放因子以提高图片质量（2x = 200% 大小）
                scale_factor = 2
                image = presentation.Slides[i].SaveAsImage(scale_factor, scale_factor)
                
                # 保存为高质量PNG
                image.Save(image_path, ImageFormat.get_Png())
                
                slide_info.image_path = image_path
                
                self.logger.debug(f"转换幻灯片 {i+1}: {image_path}")
            
            presentation.Dispose()
            return True
            
        except Exception as e:
            self.logger.error(f"Spire转换失败: {e}")
            return False
    
    def _convert_via_pdf(self) -> bool:
        """
        通过PDF中间格式转换
        
        Returns:
            bool: 是否成功
        """
        try:
            # 先转换为PDF
            pdf_path = os.path.join(self.temp_dir, "temp.pdf")
            
            # 使用LibreOffice或其他工具转换为PDF
            # 这里需要系统安装LibreOffice
            import subprocess
            
            cmd = [
                'libreoffice',
                '--headless',
                '--convert-to', 'pdf',
                '--outdir', self.temp_dir,
                str(self.ppt_path)
            ]
            
            result = subprocess.run(cmd, capture_output=True, text=True)
            if result.returncode != 0:
                self.logger.error(f"转换为PDF失败: {result.stderr}")
                return False
            
            # 从PDF转换为图片
            if os.path.exists(pdf_path):
                # 使用高DPI以获得更好的图片质量
                images = convert_from_path(
                    pdf_path, 
                    dpi=PPT_PROCESSING.get('dpi', 300),  # 使用配置的DPI，默认300
                    fmt='png',
                    thread_count=4  # 多线程加速转换
                )
                
                for i, image in enumerate(images):
                    if i < len(self.slides_info):
                        image_path = os.path.join(self.temp_dir, f"slide_{i+1:03d}.png")
                        # 保存高质量PNG
                        image.save(
                            image_path, 
                            'PNG', 
                            quality=PPT_PROCESSING.get('image_quality', 95),
                            optimize=True
                        )
                        self.slides_info[i].image_path = image_path
                        self.logger.debug(f"生成图片 {i+1}/{len(images)}: {image_path}")
                
                return True
            
            return False
            
        except Exception as e:
            self.logger.error(f"PDF转换失败: {e}")
            return False
    
    def get_slide_info(self, slide_number: int) -> Optional[SlideInfo]:
        """
        获取指定幻灯片信息
        
        Args:
            slide_number: 幻灯片编号（从1开始）
            
        Returns:
            Optional[SlideInfo]: 幻灯片信息
        """
        if 1 <= slide_number <= len(self.slides_info):
            return self.slides_info[slide_number - 1]
        return None
    
    def get_all_slides_info(self) -> List[SlideInfo]:
        """
        获取所有幻灯片信息
        
        Returns:
            List[SlideInfo]: 所有幻灯片信息
        """
        return self.slides_info.copy()
    
    def get_slide_count(self) -> int:
        """
        获取幻灯片数量
        
        Returns:
            int: 幻灯片数量
        """
        return len(self.slides_info)
    
    def get_presentation_info(self) -> Dict[str, Any]:
        """
        获取演示文稿信息
        
        Returns:
            Dict[str, Any]: 演示文稿信息
        """
        file_info = FileUtils.get_file_info(self.ppt_path)
        
        return {
            'file_path': str(self.ppt_path),
            'file_name': self.ppt_path.name,
            'file_size': file_info.get('size', 0),
            'slide_count': len(self.slides_info),
            'created_time': file_info.get('created', 0),
            'modified_time': file_info.get('modified', 0),
            'total_images': sum(slide.image_count for slide in self.slides_info),
            'total_charts': sum(slide.chart_count for slide in self.slides_info),
            'total_tables': sum(slide.table_count for slide in self.slides_info),
            'slide_types': {
                'title': sum(1 for slide in self.slides_info if slide.slide_type == 'title'),
                'section': sum(1 for slide in self.slides_info if slide.slide_type == 'section'),
                'content': sum(1 for slide in self.slides_info if slide.slide_type == 'content'),
                'summary': sum(1 for slide in self.slides_info if slide.slide_type == 'summary')
            }
        }
    
    def export_text_content(self, output_path: str) -> bool:
        """
        导出文本内容
        
        Args:
            output_path: 输出文件路径
            
        Returns:
            bool: 是否成功
        """
        try:
            content = []
            
            for slide_info in self.slides_info:
                content.append(f"# 第{slide_info.slide_number}张 - {slide_info.title}")
                content.append(f"类型: {slide_info.slide_type}")
                content.append("")
                
                if slide_info.content:
                    content.append("## 内容")
                    for text in slide_info.content:
                        content.append(f"- {text}")
                    content.append("")
                
                if slide_info.bullet_points:
                    content.append("## 要点")
                    for point in slide_info.bullet_points:
                        indent = "  " * point['level']
                        content.append(f"{indent}- {point['text']}")
                    content.append("")
                
                if slide_info.notes:
                    content.append("## 备注")
                    content.append(slide_info.notes)
                    content.append("")
                
                content.append("---")
                content.append("")
            
            return FileUtils.write_text_file(output_path, "\n".join(content))
            
        except Exception as e:
            self.logger.error(f"导出文本内容失败: {e}")
            return False
    
    def cleanup(self):
        """清理临时文件"""
        if self.temp_dir:
            FileUtils.clean_temp_dir(self.temp_dir)
            self.temp_dir = ""
    
    def __del__(self):
        """析构函数，清理资源"""
        self.cleanup()