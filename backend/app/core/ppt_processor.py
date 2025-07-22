"""
PPT处理模块 - Web版本

负责PPT文件的读取、分析和图片转换功能。
"""

import os
import time
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
from dataclasses import dataclass
import logging
import tempfile

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from pdf2image import convert_from_path
    import fitz  # PyMuPDF
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False


# 配置常量
PPT_PROCESSING = {
    'max_file_size': 100 * 1024 * 1024,  # 100MB
    'dpi': 300,
    'output_format': 'png',
    'conversion_timeout': 300,  # 5分钟
}


@dataclass
class SlideInfo:
    """幻灯片信息"""
    slide_number: int
    title: str
    content: List[str]
    bullet_points: List[Dict[str, Any]]
    image_count: int
    chart_count: int
    table_count: int
    slide_type: str
    image_path: Optional[str] = None
    notes: str = ""


class PPTProcessor:
    """PPT处理器"""
    
    def __init__(self, ppt_path: str):
        """
        初始化PPT处理器
        
        Args:
            ppt_path: PPT文件路径
        """
        self.ppt_path = Path(ppt_path)
        self.temp_dir = ""
        self.slides_info: List[SlideInfo] = []
        self.logger = logging.getLogger(__name__)
        
        # 检查文件是否存在
        if not self.ppt_path.exists():
            raise FileNotFoundError(f"PPT文件不存在: {ppt_path}")
        
        # 检查文件格式
        if not self._is_ppt_file(ppt_path):
            raise ValueError(f"不支持的文件格式: {self.ppt_path.suffix}")
        
        # 检查文件大小
        file_size = self.ppt_path.stat().st_size
        if file_size > PPT_PROCESSING['max_file_size']:
            raise ValueError(f"文件太大: {file_size / (1024*1024):.1f}MB")
        
        self.logger.info(f"初始化PPT处理器: {ppt_path}")
    
    def _is_ppt_file(self, file_path: str) -> bool:
        """检查是否是PPT文件"""
        return Path(file_path).suffix.lower() in ['.ppt', '.pptx']
    
    def process(self) -> bool:
        """
        处理PPT文件
        
        Returns:
            bool: 是否成功
        """
        try:
            # 创建临时目录
            self.temp_dir = tempfile.mkdtemp(prefix="ppt_process_")
            
            # 提取文本内容
            if not self._extract_text_content():
                return False
            
            # 转换为图片
            if not self._convert_to_images():
                return False
            
            self.logger.info(f"PPT处理完成，共{len(self.slides_info)}张幻灯片")
            return True
            
        except Exception as e:
            self.logger.error(f"PPT处理失败: {e}")
            return False
    
    def _extract_text_content(self) -> bool:
        """提取文本内容"""
        try:
            if not PPTX_AVAILABLE:
                self.logger.error("python-pptx未安装，无法提取文本内容")
                return False
            
            # 如果是.ppt文件，需要先转换为.pptx
            if self.ppt_path.suffix.lower() == '.ppt':
                self.logger.warning("暂不支持.ppt格式的文本提取，将仅进行图片转换")
                return True
            
            # 打开PPTX文件
            prs = Presentation(str(self.ppt_path))
            
            for slide_idx, slide in enumerate(prs.slides):
                # 提取标题
                title = ""
                if slide.shapes.title:
                    title = slide.shapes.title.text.strip()
                
                # 提取内容
                content = []
                bullet_points = []
                image_count = 0
                chart_count = 0
                table_count = 0
                
                for shape in slide.shapes:
                    # 文本框
                    if hasattr(shape, 'text') and shape.text.strip():
                        content.append(shape.text.strip())
                    
                    # 表格
                    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                        table_count += 1
                    
                    # 图片
                    elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image_count += 1
                    
                    # 图表
                    elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                        chart_count += 1
                
                # 提取备注
                notes = ""
                if slide.notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                
                # 判断幻灯片类型
                slide_type = self._determine_slide_type(
                    title, content, image_count, chart_count, table_count
                )
                
                # 创建幻灯片信息
                slide_info = SlideInfo(
                    slide_number=slide_idx + 1,
                    title=title,
                    content=content,
                    bullet_points=bullet_points,
                    image_count=image_count,
                    chart_count=chart_count,
                    table_count=table_count,
                    slide_type=slide_type,
                    notes=notes
                )
                
                self.slides_info.append(slide_info)
            
            self.logger.info(f"文本提取完成，共{len(self.slides_info)}张幻灯片")
            return True
            
        except Exception as e:
            self.logger.error(f"文本提取失败: {e}")
            return False
    
    def _determine_slide_type(self, title: str, content: List[str], 
                            image_count: int, chart_count: int, table_count: int) -> str:
        """判断幻灯片类型"""
        if not title and not content:
            return "empty"
        elif "目录" in title or "outline" in title.lower():
            return "toc"
        elif chart_count > 0:
            return "chart"
        elif table_count > 0:
            return "table"
        elif image_count > 0:
            return "image"
        elif len(content) > 3:
            return "content"
        else:
            return "title"
    
    def _convert_to_images(self) -> bool:
        """转换PPT为图片"""
        try:
            # 方法1: 使用LibreOffice转换
            if self._convert_with_libreoffice():
                return True
            
            # 方法2: 如果有PDF转换能力，先转PDF再转图片
            if PDF_AVAILABLE and self._convert_via_pdf():
                return True
            
            # 如果都失败了，返回False
            self.logger.error("所有图片转换方法都失败了")
            return False
            
        except Exception as e:
            self.logger.error(f"图片转换失败: {e}")
            return False
    
    def _convert_with_libreoffice(self) -> bool:
        """使用LibreOffice转换"""
        try:
            import subprocess
            
            # 检查LibreOffice是否可用
            try:
                subprocess.run(['soffice', '--version'], 
                             capture_output=True, check=True, timeout=5)
            except (subprocess.CalledProcessError, FileNotFoundError, subprocess.TimeoutExpired):
                self.logger.info("LibreOffice不可用，跳过此方法")
                return False
            
            # 创建输出目录
            output_dir = Path(self.temp_dir) / "images"
            output_dir.mkdir(exist_ok=True)
            
            # 先转换为PDF
            pdf_path = Path(self.temp_dir) / f"{self.ppt_path.stem}.pdf"
            cmd_pdf = [
                'soffice', '--headless', '--convert-to', 'pdf',
                '--outdir', str(self.temp_dir), str(self.ppt_path)
            ]
            
            result = subprocess.run(cmd_pdf, capture_output=True, text=True, 
                                  timeout=PPT_PROCESSING['conversion_timeout'])
            
            if result.returncode != 0:
                self.logger.error(f"PDF转换失败: {result.stderr}")
                return False
            
            if not pdf_path.exists():
                self.logger.error("PDF文件未生成")
                return False
            
            # 使用pdf2image转换PDF为图片
            if not PDF_AVAILABLE:
                self.logger.error("pdf2image不可用")
                return False
            
            images = convert_from_path(
                str(pdf_path),
                dpi=PPT_PROCESSING['dpi'],
                output_folder=str(output_dir),
                fmt='png'
            )
            
            # 更新幻灯片信息中的图片路径
            for i, image in enumerate(images):
                if i < len(self.slides_info):
                    image_path = output_dir / f"slide_{i+1:03d}.png"
                    image.save(image_path)
                    self.slides_info[i].image_path = str(image_path)
            
            self.logger.info(f"LibreOffice转换完成，生成{len(images)}张图片")
            return True
            
        except Exception as e:
            self.logger.error(f"LibreOffice转换失败: {e}")
            return False
    
    def _convert_via_pdf(self) -> bool:
        """通过PDF转换"""
        # 这里可以实现其他PDF转换方法
        self.logger.info("PDF转换方法暂未实现")
        return False
    
    def get_slides_info(self) -> List[SlideInfo]:
        """获取幻灯片信息"""
        return self.slides_info
    
    def get_slide_image_paths(self) -> List[str]:
        """获取幻灯片图片路径"""
        paths = []
        for slide in self.slides_info:
            if slide.image_path:
                paths.append(slide.image_path)
        return paths
    
    def get_total_slides(self) -> int:
        """获取幻灯片总数"""
        return len(self.slides_info)
    
    def cleanup(self):
        """清理临时文件"""
        if self.temp_dir and Path(self.temp_dir).exists():
            try:
                import shutil
                shutil.rmtree(self.temp_dir)
                self.logger.info(f"清理临时目录: {self.temp_dir}")
            except Exception as e:
                self.logger.error(f"清理临时目录失败: {e}")
    
    def __del__(self):
        """析构函数"""
        self.cleanup()